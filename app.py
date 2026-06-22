"""
Impulse Academy – Générateur d'attestations de formation
Backend Flask – Version avec authentification et formations prédéfinies
"""
import os, re, uuid, zipfile, subprocess, shutil, hashlib
from pathlib import Path
from functools import wraps
from flask import (Flask, request, jsonify, send_file,
                   render_template, send_from_directory,
                   session, redirect, url_for)
from pptx import Presentation
import pptx.oxml.ns as ns_pptx
import openpyxl

app = Flask(__name__)
app.secret_key = "ImpulseAcademy_S3cr3t_2026!"
app.config['MAX_CONTENT_LENGTH'] = 16 * 1024 * 1024

BASE_DIR      = Path(__file__).parent
MODELE_PPTX   = BASE_DIR / "modele_global.pptx"
GENERATED_DIR = BASE_DIR / "generated"
UPLOADS_DIR   = BASE_DIR / "uploads"
GENERATED_DIR.mkdir(exist_ok=True)
UPLOADS_DIR.mkdir(exist_ok=True)

# ── AUTHENTIFICATION ───────────────────────────────────────────
USERS = {
    "aymane.akdahi@gmail.com":      hashlib.sha256(b"Impulse2026@cademy!").hexdigest(),
    "tarik.beladgham@impulse-aca.fr": hashlib.sha256(b"Impulse2026@cademy!").hexdigest(),
    "khalil.kriach@impulse-aca.fr":  hashlib.sha256(b"Impulse2026@cademy!").hexdigest(),
}

def login_required(f):
    @wraps(f)
    def decorated(*args, **kwargs):
        if not session.get('logged_in'):
            if request.is_json or request.path.startswith('/api/'):
                return jsonify({'error': 'Non authentifié'}), 401
            return redirect(url_for('login_page'))
        return f(*args, **kwargs)
    return decorated

@app.route('/login', methods=['GET'])
def login_page():
    if session.get('logged_in'):
        return redirect(url_for('index'))
    return render_template('login.html')

@app.route('/api/login', methods=['POST'])
def api_login():
    data  = request.get_json() or {}
    email = (data.get('email') or '').strip().lower()
    pwd   = (data.get('password') or '').strip()
    h     = hashlib.sha256(pwd.encode()).hexdigest()
    if email in USERS and USERS[email] == h:
        session['logged_in'] = True
        session['email']     = email
        return jsonify({'ok': True})
    return jsonify({'error': 'Email ou mot de passe incorrect'}), 401

@app.route('/api/logout', methods=['POST'])
def api_logout():
    session.clear()
    return jsonify({'ok': True})

# ── FORMATIONS PRÉDÉFINIES ─────────────────────────────────────
# rpps_required : True = RPPS obligatoire, False = ligne RPPS supprimée
FORMATIONS = [
    {
        "id": "f1",
        "label": "Formation rayonnements ionisants – Module spécifique \"Infirmier de santé au travail\" – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "rayonnements ionisants",
        "fonction": "Infirmier de santé au travail",
        "rpps_required": False,
    },
    {
        "id": "f2",
        "label": "Formation rayonnements ionisants – Module spécifique \"Médecin du travail\" – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "rayonnements ionisants",
        "fonction": "Médecin du travail",
        "rpps_required": True,
    },
    {
        "id": "f3",
        "label": "Formation rayonnements ionisants – Recyclage module spécifique \"Infirmier de santé au travail\" – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "rayonnements ionisants",
        "fonction": "Recyclage module spécifique – Infirmier de santé au travail",
        "rpps_required": False,
    },
    {
        "id": "f4",
        "label": "Formation rayonnements ionisants – Recyclage module spécifique \"Médecin du travail\" – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "rayonnements ionisants",
        "fonction": "Recyclage module spécifique – Médecin du travail",
        "rpps_required": True,
    },
    {
        "id": "f5",
        "label": "Formation aux rayonnements ionisants – Module complémentaire B : Suivi des travailleurs exposés au radon – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "aux rayonnements ionisants",
        "fonction": "Module complémentaire B – Suivi des travailleurs exposés au radon",
        "rpps_required": True,
    },
    {
        "id": "f6",
        "label": "Formation aux rayonnements ionisants – Module complémentaire C : Suivi des travailleurs intervenant en situation d'urgence radiologique – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "aux rayonnements ionisants",
        "fonction": "Module complémentaire C – Suivi des travailleurs intervenant en situation d'urgence radiologique",
        "rpps_required": True,
    },
    {
        "id": "f7",
        "label": "Formation aux rayonnements ionisants – Module complémentaire D : Suivi des travailleurs exposés aux neutrons – Conforme à l'arrêté du 06 août 2024",
        "type_formation": "aux rayonnements ionisants",
        "fonction": "Module complémentaire D – Suivi des travailleurs exposés aux neutrons",
        "rpps_required": True,
    },
]

FORMATIONS_BY_ID = {f["id"]: f for f in FORMATIONS}

# ── HELPERS PPTX ──────────────────────────────────────────────

def supprimer_highlight(run):
    rPr = run._r.find(ns_pptx.qn('a:rPr'))
    if rPr is not None:
        hl = rPr.find(ns_pptx.qn('a:highlight'))
        if hl is not None:
            rPr.remove(hl)

def supprimer_highlight_shape(shape):
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        for run in para.runs:
            supprimer_highlight(run)

def remplacer_placeholder_runs(shape, placeholder, valeur):
    """
    Remplace un placeholder {{...}} qui peut être fragmenté sur plusieurs runs.
    Stratégie : reconstituer le texte du para, trouver le placeholder,
    puis mettre toute la valeur dans le premier run qui contient le début
    du placeholder et vider les autres runs du placeholder.
    """
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        runs = para.runs
        # Construire le texte complet du para avec indices de runs
        texte = "".join(r.text or "" for r in runs)
        if placeholder not in texte: continue

        # Reconstruire: mettre la valeur dans le premier run qui contribue
        # au placeholder, vider tous les runs intermédiaires du placeholder
        new_texte = texte.replace(placeholder, valeur)
        # Stratégie simple: concentrer tout dans run[0], vider les autres
        if runs:
            runs[0].text = new_texte
            supprimer_highlight(runs[0])
            for r in runs[1:]:
                r.text = ""

def supprimer_paragraphe_par_placeholder(shape, placeholder):
    """Supprime le paragraphe XML entier contenant le placeholder."""
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        texte = "".join(r.text or "" for r in para.runs)
        if placeholder in texte:
            p_elem = para._p
            p_elem.getparent().remove(p_elem)

def vider_placeholder_optionnel(shape, placeholder_tag, valeur=None):
    """
    Gère les blocs Optionnelle_ :
    - Si valeur fournie : remplace {{Optionnelle_}} par '' et garde le reste
    - Si valeur None (absent) : supprime le paragraphe entier
    """
    if not shape.has_text_frame: return
    for para in shape.text_frame.paragraphs:
        texte = "".join(r.text or "" for r in para.runs)
        if placeholder_tag not in texte: continue
        if valeur is None:
            # Supprimer le paragraphe entier
            para._p.getparent().remove(para._p)
        else:
            # Supprimer juste le marqueur optionnel, garder le reste
            runs = para.runs
            new_texte = texte.replace(placeholder_tag, "")
            if runs:
                runs[0].text = new_texte
                supprimer_highlight(runs[0])
                for r in runs[1:]:
                    r.text = ""

def _re_placeholder(nom):
    """Regex qui matche {{ NOM }} avec ou sans espaces intérieurs."""
    return re.compile(r'\{+\s*\{+\s*' + re.escape(nom) + r'\s*\}+\s*\}+')

def _consolider_para(para, new_text):
    """Met new_text dans run[0] du para, vide les autres runs."""
    runs = para.runs
    if not runs: return
    runs[0].text = new_text
    supprimer_highlight(runs[0])
    for r in runs[1:]:
        r.text = ""

def _texte_para(para):
    return "".join(r.text or "" for r in para.runs)

def generer_pptx(participant, session_data, formation, pptx_sortie):
    """
    Génère un PPTX à partir du modèle global avec les placeholders {{...}}.
    Tolère les espaces variables à l'intérieur des {{ }}.
    """
    prs   = Presentation(str(MODELE_PPTX))
    slide = prs.slides[0]

    nom        = participant['nom']
    rpps       = participant.get('rpps', '').strip()
    date_debut = session_data['dateDebut']
    date_fin   = session_data['dateFin']
    date_deliv = session_data['dateDeliv']
    duree      = session_data.get('dureeValidite', '5 ans')
    rpps_requis = formation['rpps_required']

    # Construire la chaîne de dates de formation
    if date_debut == date_fin:
        dates_form = date_debut
    else:
        d1 = date_debut.split('/')
        d2 = date_fin.split('/')
        if len(d1) == 3 and len(d2) == 3 and d1[1] == d2[1] and d1[2] == d2[2]:
            dates_form = f"{d1[0]} au {date_fin}"
        else:
            dates_form = f"{date_debut} au {date_fin}"

    # Regex pour tous les placeholders connus
    RE_NOM      = _re_placeholder('NOM_PARTICIPANT')
    RE_OPT      = re.compile(r'\{+\s*\{+\s*Optionnelle_\s*\}+\s*\}+')
    RE_RPPS_NUM = re.compile(r'\{+\s*\{+\s*Optionnelle_Numero_RPPS\s*\}+\s*\}+')
    RE_TYPE     = _re_placeholder('TYPE_FORMATION')
    RE_FONC     = _re_placeholder('FONCTION')
    RE_DATES    = _re_placeholder('DATES_FORMATION')
    RE_DELIV    = _re_placeholder('DATE_DELIVRANCE')
    RE_DUREE    = re.compile(r'\{+\s*\{+\s*Optionnelle_DUREE_VALIDITE\s*\}+\s*\}+')
    # Regex générique pour nettoyer tout placeholder résiduel
    RE_ANY_PH   = re.compile(r'\{+[^}]+\}+')

    for shape in slide.shapes:
        if not shape.has_text_frame: continue

        # Supprimer tous les highlights
        supprimer_highlight_shape(shape)

        texte_shape = "".join(_texte_para(p) for p in shape.text_frame.paragraphs)

        # ── Shape 73 : NOM + RPPS ──
        if RE_NOM.search(texte_shape):
            for para in shape.text_frame.paragraphs:
                t = _texte_para(para)
                if RE_NOM.search(t):
                    _consolider_para(para, RE_NOM.sub(nom, t))

        # Traitement du bloc RPPS (après remplacement du nom)
        texte_shape = "".join(_texte_para(p) for p in shape.text_frame.paragraphs)
        if 'Optionnelle_Numero_RPPS' in texte_shape or ('Optionnelle_' in texte_shape and 'RPPS' in texte_shape.upper()):
            if rpps_requis and rpps:
                # Garder la ligne RPPS : remplacer le bloc entier par "N°RPPS: <numero>"
                # Le modèle contient : " {{Optionnelle_}} : N°RPPS: {{Optionnelle_Numero_RPPS}}"
                # On reconstruit proprement
                for para in shape.text_frame.paragraphs:
                    t = _texte_para(para)
                    if 'Optionnelle_Numero_RPPS' in t or ('Optionnelle_' in t and 'RPPS' in t.upper()):
                        # Remplacer tout le bloc optionnel+RPPS par la valeur propre
                        new_t = re.sub(
                            r'\s*\{+\s*\{+\s*Optionnelle_\s*\}+\s*\}+\s*[:\-]?\s*N°RPPS\s*[:\-]?\s*\{+\s*\{+\s*Optionnelle_Numero_RPPS\s*\}+\s*\}+',
                            f'N°RPPS : {rpps}',
                            t, flags=re.IGNORECASE
                        )
                        # Si la regex n'a rien changé, nettoyage manuel
                        if new_t == t:
                            new_t = RE_OPT.sub('', t)
                            new_t = RE_RPPS_NUM.sub(rpps, new_t)
                            # Supprimer le ":" parasite au début
                            new_t = re.sub(r'^\s*[:\-]\s*', '', new_t.strip())
                        new_t = new_t.strip()
                        if new_t:
                            _consolider_para(para, new_t)
                        else:
                            try: para._p.getparent().remove(para._p)
                            except: pass
            else:
                # Supprimer le paragraphe contenant RPPS entièrement
                paras_rpps = [
                    para._p for para in shape.text_frame.paragraphs
                    if 'RPPS' in _texte_para(para).upper() or 'Optionnelle_' in _texte_para(para)
                ]
                for p_elem in paras_rpps:
                    try: p_elem.getparent().remove(p_elem)
                    except: pass

        # ── Shape 74 : TITRE (TYPE_FORMATION + FONCTION) ──
        texte_shape = "".join(_texte_para(p) for p in shape.text_frame.paragraphs)
        if RE_TYPE.search(texte_shape) or RE_FONC.search(texte_shape):
            for para in shape.text_frame.paragraphs:
                t = _texte_para(para)
                if RE_TYPE.search(t) or RE_FONC.search(t):
                    new_t = RE_TYPE.sub(formation['type_formation'], t)
                    new_t = RE_FONC.sub(formation['fonction'], new_t)
                    _consolider_para(para, new_t)

        # ── Shape 75 : DATES DE FORMATION ──
        texte_shape = "".join(_texte_para(p) for p in shape.text_frame.paragraphs)
        if RE_DATES.search(texte_shape):
            for para in shape.text_frame.paragraphs:
                t = _texte_para(para)
                if RE_DATES.search(t):
                    _consolider_para(para, RE_DATES.sub(dates_form, t))

        # ── Shape 76 : DATE DÉLIVRANCE + DURÉE ──
        texte_shape = "".join(_texte_para(p) for p in shape.text_frame.paragraphs)
        if RE_DELIV.search(texte_shape):
            for para in shape.text_frame.paragraphs:
                t = _texte_para(para)
                if RE_DELIV.search(t):
                    new_t = RE_DELIV.sub(date_deliv, t)
                    if not duree:
                        # Supprimer la partie durée : tout depuis le \t jusqu'à la fin
                        new_t = re.sub(r'\s*\t\s*.*$', '', new_t, flags=re.DOTALL)
                        # Nettoyer les marqueurs optionnels résiduels
                        new_t = RE_OPT.sub('', new_t)
                        new_t = RE_DUREE.sub('', new_t)
                        new_t = re.sub(r'\s*[:\-]\s*$', '', new_t.rstrip())
                    else:
                        # Remplacer le bloc optionnel+durée entier proprement
                        # Modèle : "\t {{Optionnelle_}}  :      Durée de validité: {{Optionnelle_DUREE_VALIDITE}}"
                        new_t = re.sub(
                            r'\t\s*\{+\s*\{+\s*Optionnelle_\s*\}+\s*\}+\s*[:\-]?\s*Durée de validité\s*[:\-]\s*\{+\s*\{+\s*Optionnelle_DUREE_VALIDITE\s*\}+\s*\}+',
                            f'\tDurée de validité : {duree}',
                            new_t, flags=re.IGNORECASE
                        )
                        # Si pas changé, nettoyage manuel
                        new_t = RE_OPT.sub('', new_t)
                        new_t = RE_DUREE.sub(duree, new_t)
                        # Supprimer le ":" parasite restant avant "Durée"
                        new_t = re.sub(r'\t\s*[:\-]\s*', '\t', new_t)
                    _consolider_para(para, new_t.rstrip())

        # Nettoyage final : supprimer tout placeholder {{...}} résiduel
        for para in shape.text_frame.paragraphs:
            t = _texte_para(para)
            if RE_ANY_PH.search(t):
                _consolider_para(para, RE_ANY_PH.sub('', t))

    prs.save(str(pptx_sortie))


LIBREOFFICE_AVAILABLE = shutil.which("libreoffice") is not None

def pptx_vers_pdf(pptx_path, output_dir):
    if not LIBREOFFICE_AVAILABLE:
        return None  # Pas de LibreOffice → PDF non disponible
    res = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf",
         "--outdir", str(output_dir), str(pptx_path)],
        capture_output=True, text=True, timeout=120)
    if res.returncode != 0:
        raise RuntimeError(res.stderr[:300])
    pdf = Path(output_dir) / f"{Path(pptx_path).stem}.pdf"
    if not pdf.exists():
        raise FileNotFoundError(str(pdf))
    return pdf


def sanitize(name):
    return re.sub(r'[^\w\-]', '_', name).strip('_')

# ── ROUTES ────────────────────────────────────────────────────

@app.route('/')
@login_required
def index():
    return render_template('index.html')

@app.route('/api/formations')
@login_required
def api_formations():
    return jsonify(FORMATIONS)

@app.route('/api/upload-modele', methods=['POST'])
@login_required
def upload_modele():
    if 'file' not in request.files:
        return jsonify({'error': 'Aucun fichier'}), 400
    file = request.files['file']
    if not file.filename.lower().endswith('.pptx'):
        return jsonify({'error': 'Le fichier doit être un .pptx'}), 400
    file.save(str(MODELE_PPTX))
    return jsonify({'ok': True, 'message': 'Modèle enregistré'})

@app.route('/api/import-excel', methods=['POST'])
@login_required
def import_excel():
    if 'file' not in request.files:
        return jsonify({'error': 'Aucun fichier'}), 400
    file = request.files['file']
    fname = file.filename.lower()
    participants = []
    try:
        if fname.endswith('.csv') or fname.endswith('.txt'):
            content = file.read().decode('utf-8-sig')
            lines = content.split('\n')
            header = None
            for i, line in enumerate(lines):
                line = line.strip()
                if not line: continue
                cols = [c.strip() for c in re.split(r'[;\t,]', line)]
                if i == 0 and any(k in c.lower() for c in cols for k in ['nom','prenom','name','rpps']):
                    header = [c.lower().replace(' ', '_') for c in cols]
                    continue
                if header:
                    d = dict(zip(header, cols))
                    nom  = d.get('nom_prenom') or d.get('nom') or (cols[0] if cols else '')
                    rpps = d.get('rpps', '')
                else:
                    nom  = cols[0] if cols else ''
                    rpps = cols[1] if len(cols) > 1 else ''
                if nom:
                    participants.append({'nom': nom, 'rpps': rpps})
        elif fname.endswith('.xlsx') or fname.endswith('.xls'):
            tmp = UPLOADS_DIR / f"{uuid.uuid4()}.xlsx"
            file.save(str(tmp))
            wb = openpyxl.load_workbook(str(tmp), read_only=True)
            ws = wb.active
            header = None
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if not any(row): continue
                vals = [str(c).strip() if c else '' for c in row]
                if i == 0 and any(k in v.lower() for v in vals for k in ['nom','prenom','name','rpps']):
                    header = [v.lower().replace(' ', '_') for v in vals]
                    continue
                if header:
                    d = dict(zip(header, vals))
                    nom  = d.get('nom_prenom') or d.get('nom') or vals[0]
                    rpps = d.get('rpps', '')
                else:
                    nom  = vals[0]
                    rpps = vals[1] if len(vals) > 1 else ''
                if nom and nom not in ('None', ''):
                    participants.append({'nom': nom, 'rpps': rpps})
            wb.close()
            tmp.unlink(missing_ok=True)
        else:
            return jsonify({'error': 'Format non supporté (.xlsx, .csv)'}), 400
    except Exception as e:
        return jsonify({'error': str(e)}), 500
    return jsonify({'participants': participants, 'count': len(participants)})

@app.route('/api/generer', methods=['POST'])
@login_required
def generer():
    data = request.get_json()
    if not data:
        return jsonify({'error': 'Données manquantes'}), 400

    formation_id = data.get('formationId')
    if not formation_id or formation_id not in FORMATIONS_BY_ID:
        return jsonify({'error': 'Formation invalide'}), 400

    formation = FORMATIONS_BY_ID[formation_id]
    rpps_requis = formation['rpps_required']

    for field in ['dateDebut', 'dateFin', 'dateDeliv', 'participants']:
        if not data.get(field):
            return jsonify({'error': f'Champ manquant: {field}'}), 400

    participants = data['participants']
    if not participants:
        return jsonify({'error': 'Aucun participant'}), 400

    # Vérification RPPS obligatoire
    if rpps_requis:
        manquants = [p['nom'] for p in participants if not p.get('rpps', '').strip()]
        if manquants:
            return jsonify({
                'error': f"RPPS manquant pour : {', '.join(manquants)}"
            }), 400

    session_data = {
        'dateDebut':    data.get('dateDebut', ''),
        'dateFin':      data.get('dateFin', ''),
        'dateDeliv':    data.get('dateDeliv', ''),
        'dureeValidite': data.get('dureeValidite', '5 ans'),
    }

    job_id   = str(uuid.uuid4())[:8]
    job_dir  = GENERATED_DIR / job_id
    pptx_dir = job_dir / 'pptx'
    pdf_dir  = job_dir / 'pdf'
    pptx_dir.mkdir(parents=True)
    pdf_dir.mkdir(parents=True)

    generated, errors = [], []

    for p in participants:
        nom  = p.get('nom', '').strip()
        rpps = p.get('rpps', '').strip() if rpps_requis else ''
        if not nom: continue
        safe = sanitize(nom)
        pptx_path = pptx_dir / f"Attestation_{safe}.pptx"
        try:
            generer_pptx({'nom': nom, 'rpps': rpps}, session_data, formation, pptx_path)
            pdf_path = pptx_vers_pdf(pptx_path, pdf_dir)
            entry = {
                'nom':  nom,
                'pptx': pptx_path.name,
            }
            if pdf_path:
                entry['pdf_file'] = pdf_path.name
                entry['pdf_size'] = pdf_path.stat().st_size
            else:
                entry['pdf_file'] = ''
                entry['pdf_size'] = 0
            generated.append(entry)
        except Exception as e:
            errors.append({'nom': nom, 'erreur': str(e)})

    if not generated:
        shutil.rmtree(job_dir, ignore_errors=True)
        msg = errors[0]['erreur'] if errors else 'Erreur inconnue'
        return jsonify({'error': msg}), 500

    safe_fn  = sanitize(formation['id'])
    zip_all  = job_dir / f"Attestations_{safe_fn}.zip"
    zip_pdf  = job_dir / f"PDF_{safe_fn}.zip"

    has_pdf = any(list(pdf_dir.glob('*.pdf')))

    with zipfile.ZipFile(str(zip_all), 'w', zipfile.ZIP_DEFLATED) as zf:
        for f in pdf_dir.glob('*.pdf'):   zf.write(f, f"PDF/{f.name}")
        for f in pptx_dir.glob('*.pptx'): zf.write(f, f"PPTX/{f.name}")

    if has_pdf:
        with zipfile.ZipFile(str(zip_pdf), 'w', zipfile.ZIP_DEFLATED) as zf:
            for f in pdf_dir.glob('*.pdf'): zf.write(f, f.name)

    return jsonify({
        'success':   True,
        'count':     len(generated),
        'errors':    errors,
        'generated': generated,
        'job_id':    job_id,
        'has_pdf':   has_pdf,
        'zip_url':   f'/download/{job_id}/zip',
        'pdf_url':   f'/download/{job_id}/pdf' if has_pdf else None,
    })

@app.route('/download/<job_id>/zip')
@login_required
def dl_zip(job_id):
    zips = list((GENERATED_DIR / job_id).glob('Attestations_*.zip'))
    if not zips: return "Non trouvé", 404
    return send_file(str(zips[0]), as_attachment=True, download_name=zips[0].name)

@app.route('/download/<job_id>/pdf')
@login_required
def dl_pdf(job_id):
    zips = list((GENERATED_DIR / job_id).glob('PDF_*.zip'))
    if not zips: return "Non trouvé", 404
    return send_file(str(zips[0]), as_attachment=True, download_name=zips[0].name)

if __name__ == '__main__':
    print("Impulse Academy – http://localhost:5000")
    app.run(debug=False, host='0.0.0.0', port=5000)
